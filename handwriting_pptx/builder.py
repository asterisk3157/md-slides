"""pptx (zip) パッケージング。

**戦略**: python-pptx が生成する pptx の骨格（slideMaster / slideLayout /
theme / presentation.xml など）は PowerPoint 互換が確実なので、これを
ベースとして使い、自前のスライドXML（grpSp + contentPart + timing）
と ink ファイル + media ファイル + 関連メタを後から注入する。

自前で書いていた SLIDE_MASTER / SLIDE_LAYOUT / THEME 等の静的テンプレートは
要素不足で PowerPoint が拒否したため廃止。
"""
from __future__ import annotations

import io
import os
import tempfile
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import List

from PIL import Image as PILImage
from pptx import Presentation as PPTXPresentation
from pptx.util import Emu

from .units import cm_to_emu, DEFAULT_SLIDE_W_CM, DEFAULT_SLIDE_H_CM


# ---------- パッケージング用 DTO ----------

@dataclass
class SlideData:
    """1スライドぶんのデータ。"""

    slide_xml: str                       # 完成済み (timing含む) のスライドXML
    ink_xmls: List[str] = field(default_factory=list)        # 各 inkN.xml の内容
    fallback_png_bytes: List[bytes] = field(default_factory=list)  # 各 imageN.png のバイト列


def transparent_png_bytes(size: int = 4) -> bytes:
    """透明 PNG (size x size) のバイト列を返す。"""
    img = PILImage.new("RGBA", (size, size), (0, 0, 0, 0))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def build_pptx(slides: List[SlideData], out_path: str | Path,
               slide_w_cm: float = DEFAULT_SLIDE_W_CM,
               slide_h_cm: float = DEFAULT_SLIDE_H_CM) -> None:
    """pptx を zip としてビルドして書き出す。

    python-pptx で生成した空の pptx 骨格をベースに、自前のスライドXML
    （contentPart + timing）と ink ファイル + media ファイルを注入する。

    python-pptx の骨格は PowerPoint 互換が確実な slideMaster / slideLayout /
    theme / presentation.xml を含むため、これに乗っかるのが最も安全。
    """
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    # ---- ステップ1: python-pptx でベース pptx を tempfile に生成 ----
    pp = PPTXPresentation()
    pp.slide_width = Emu(cm_to_emu(slide_w_cm))
    pp.slide_height = Emu(cm_to_emu(slide_h_cm))
    blank_layout = pp.slide_layouts[6]  # 完全に白紙のレイアウト
    for _ in slides:
        pp.slides.add_slide(blank_layout)
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        base_path = f.name
    try:
        pp.save(base_path)

        # ---- ステップ2: 注入してzip再生成 ----
        with zipfile.ZipFile(base_path) as src:
            with zipfile.ZipFile(out_path, "w", compression=zipfile.ZIP_DEFLATED) as dst:
                ink_counter = 0
                img_counter = 0

                for name in src.namelist():
                    data = src.read(name)

                    if (name.startswith("ppt/slides/slide")
                            and name.endswith(".xml")
                            and "/_rels/" not in name):
                        idx = int(name[len("ppt/slides/slide"):-len(".xml")])
                        if 1 <= idx <= len(slides):
                            data = slides[idx - 1].slide_xml.encode("utf-8")

                    elif (name.startswith("ppt/slides/_rels/slide")
                          and name.endswith(".xml.rels")):
                        idx = int(name[len("ppt/slides/_rels/slide"):-len(".xml.rels")])
                        if 1 <= idx <= len(slides):
                            sd = slides[idx - 1]
                            existing = data.decode("utf-8")
                            add_rels: List[str] = []
                            for i, _ in enumerate(sd.ink_xmls):
                                ink_counter += 1
                                rid = sd.__dict__["_ink_rel_ids"][i]
                                add_rels.append(
                                    f'<Relationship Id="{rid}" '
                                    f'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/customXml" '
                                    f'Target="../ink/ink{ink_counter}.xml"/>'
                                )
                            for i, _ in enumerate(sd.fallback_png_bytes):
                                img_counter += 1
                                rid = sd.__dict__["_img_rel_ids"][i]
                                add_rels.append(
                                    f'<Relationship Id="{rid}" '
                                    f'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" '
                                    f'Target="../media/image{img_counter}.png"/>'
                                )
                            data = existing.replace(
                                "</Relationships>",
                                "".join(add_rels) + "</Relationships>",
                            ).encode("utf-8")

                    elif name == "[Content_Types].xml":
                        existing = data.decode("utf-8")
                        # 追加要素を </Types> の直前に挿入（順序不問・安全）
                        add_parts: List[str] = []
                        if 'Extension="png"' not in existing:
                            add_parts.append('<Default Extension="png" ContentType="image/png"/>')
                        n = 0
                        for sd in slides:
                            for _ in sd.ink_xmls:
                                n += 1
                                add_parts.append(
                                    f'<Override PartName="/ppt/ink/ink{n}.xml" '
                                    f'ContentType="application/inkml+xml"/>'
                                )
                        data = existing.replace(
                            "</Types>",
                            "".join(add_parts) + "</Types>",
                        ).encode("utf-8")

                    dst.writestr(name, data)

                # ---- ステップ3: ink + media ファイルを追加 ----
                ink_counter = 0
                img_counter = 0
                for sd in slides:
                    for ink_xml in sd.ink_xmls:
                        ink_counter += 1
                        dst.writestr(f"ppt/ink/ink{ink_counter}.xml", ink_xml)
                    for png in sd.fallback_png_bytes:
                        img_counter += 1
                        dst.writestr(f"ppt/media/image{img_counter}.png", png)
    finally:
        try:
            os.unlink(base_path)
        except OSError:
            pass
